#!/usr/bin/env python3
from __future__ import annotations

import argparse
import datetime as dt
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


@dataclass(frozen=True)
class Competitor:
    name: str
    positioning: str
    strengths: list[str]
    notable_gaps: list[str]


def _rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


THEME = {
    "navy": _rgb("#0B1F3B"),
    "blue": _rgb("#2563EB"),
    "gray": _rgb("#475569"),
    "light_gray": _rgb("#E2E8F0"),
    "black": _rgb("#0F172A"),
}


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    w, h = prs.slide_width, prs.slide_height

    header = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE (avoid extra import)
        0,
        0,
        w,
        int(h * 0.25),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME["navy"]
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), w - Inches(1.6), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(38)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    subtitle_box = slide.shapes.add_textbox(Inches(0.85), Inches(2.25), w - Inches(1.7), Inches(1.2))
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(18)
    run2.font.color.rgb = THEME["gray"]

    footer = slide.shapes.add_textbox(Inches(0.8), h - Inches(0.6), w - Inches(1.6), Inches(0.4))
    ftf = footer.text_frame
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.RIGHT
    fr = fp.add_run()
    fr.text = dt.date.today().strftime("%B %-d, %Y") if hasattr(dt.date.today(), "strftime") else str(dt.date.today())
    fr.font.size = Pt(12)
    fr.font.color.rgb = THEME["gray"]


def add_bullets_slide(prs: Presentation, title: str, bullets: Iterable[str], note: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, prs, title)

    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), prs.slide_width - Inches(1.8), prs.slide_height - Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = b
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = THEME["black"]

    if note:
        slide.notes_slide.notes_text_frame.text = note


def _add_header(slide, prs: Presentation, title: str) -> None:
    w = prs.slide_width
    header_h = Inches(0.75)
    header = slide.shapes.add_shape(1, 0, 0, w, header_h)
    header.fill.solid()
    header.fill.fore_color.rgb = THEME["navy"]
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.17), w - Inches(1.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def add_competitor_snapshot_slide(prs: Presentation, competitors: list[Competitor]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, prs, "Top competitors (snapshot)")

    cols = 3
    left = Inches(0.65)
    top = Inches(1.25)
    gap = Inches(0.3)
    total_w = prs.slide_width - Inches(1.3)
    card_w = (total_w - gap * (cols - 1)) / cols
    card_h = Inches(5.6)

    for idx, comp in enumerate(competitors[:3]):
        x = left + idx * (card_w + gap)
        card = slide.shapes.add_shape(1, x, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = THEME["light_gray"]
        card.line.width = Pt(1)

        name_box = slide.shapes.add_textbox(x + Inches(0.25), top + Inches(0.25), card_w - Inches(0.5), Inches(0.5))
        ntf = name_box.text_frame
        np = ntf.paragraphs[0]
        nr = np.add_run()
        nr.text = comp.name
        nr.font.size = Pt(22)
        nr.font.bold = True
        nr.font.color.rgb = THEME["black"]

        pos_box = slide.shapes.add_textbox(x + Inches(0.25), top + Inches(0.85), card_w - Inches(0.5), Inches(0.9))
        ptf = pos_box.text_frame
        pp = ptf.paragraphs[0]
        pr = pp.add_run()
        pr.text = comp.positioning
        pr.font.size = Pt(14)
        pr.font.color.rgb = THEME["gray"]

        strengths_box = slide.shapes.add_textbox(x + Inches(0.25), top + Inches(1.65), card_w - Inches(0.5), Inches(2.35))
        stf = strengths_box.text_frame
        stf.word_wrap = True
        stf.clear()
        p0 = stf.paragraphs[0]
        p0.text = "Strengths"
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = THEME["blue"]

        for s in comp.strengths:
            sp = stf.add_paragraph()
            sp.text = s
            sp.font.size = Pt(13)
            sp.font.color.rgb = THEME["black"]

        gaps_box = slide.shapes.add_textbox(x + Inches(0.25), top + Inches(4.15), card_w - Inches(0.5), Inches(1.2))
        gtf = gaps_box.text_frame
        gtf.word_wrap = True
        gtf.clear()
        gp0 = gtf.paragraphs[0]
        gp0.text = "Watch-outs"
        gp0.font.size = Pt(14)
        gp0.font.bold = True
        gp0.font.color.rgb = THEME["blue"]

        for g in comp.notable_gaps:
            gp = gtf.add_paragraph()
            gp.text = g
            gp.font.size = Pt(13)
            gp.font.color.rgb = THEME["black"]


def add_matrix_slide(
    prs: Presentation,
    title: str,
    rows: list[tuple[str, list[str]]],
    col_names: list[str],
    note: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, prs, title)

    table_left = Inches(0.6)
    table_top = Inches(1.35)
    table_width = prs.slide_width - Inches(1.2)
    table_height = prs.slide_height - Inches(2.1)
    table = slide.shapes.add_table(len(rows) + 1, len(col_names) + 1, table_left, table_top, table_width, table_height).table

    table.columns[0].width = Inches(4.5)
    remaining_emu = int(table_width) - int(table.columns[0].width)
    each_emu = remaining_emu // len(col_names)
    for i in range(1, len(col_names) + 1):
        table.columns[i].width = each_emu
    used_emu = each_emu * len(col_names)
    if used_emu < remaining_emu:
        table.columns[len(col_names)].width = each_emu + (remaining_emu - used_emu)

    hdr0 = table.cell(0, 0)
    hdr0.text = "Capability"
    _style_cell(hdr0, bold=True, fill=THEME["light_gray"])

    for idx, name in enumerate(col_names, start=1):
        cell = table.cell(0, idx)
        cell.text = name
        _style_cell(cell, bold=True, fill=THEME["light_gray"])

    for r, (cap, vals) in enumerate(rows, start=1):
        c0 = table.cell(r, 0)
        c0.text = cap
        _style_cell(c0)
        for c, v in enumerate(vals, start=1):
            cell = table.cell(r, c)
            cell.text = v
            _style_cell(cell, align="center")

    if note:
        slide.notes_slide.notes_text_frame.text = note


def _style_cell(cell, *, bold: bool = False, fill: RGBColor | None = None, align: str = "left") -> None:
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(12.5)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = THEME["black"]
        if align == "center":
            paragraph.alignment = PP_ALIGN.CENTER


def build_deck(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    competitors = [
        Competitor(
            name="Meltwater",
            positioning="Enterprise-grade breadth: news + social + analytics depth",
            strengths=[
                "Broadest multi-channel coverage (incl. mobile app)",
                "Strong dashboards + BI/export APIs",
                "LLM monitoring (“GenAI Lens”) + innovation",
            ],
            notable_gaps=[
                "Complex UX; learning curve",
                "Renewal price increases common",
            ],
        ),
        Competitor(
            name="Cision",
            positioning="Legacy PR standard: monitoring + database + PR Newswire workflow",
            strengths=[
                "Largest media database + PR workflows",
                "React Score for crisis/risk detection",
                "PR Newswire closed-loop distribution → pickup tracking",
            ],
            notable_gaps=[
                "Complex platform (many acquisitions stitched)",
                "Boolean building still hard for most users",
            ],
        ),
        Competitor(
            name="Muck Rack",
            positioning="Modern UX-first PR platform: fast monitoring + pitching + reporting",
            strengths=[
                "Best-in-class usability; fast onboarding",
                "Unified monitoring across partners (Nexis/TVEyes/Keyhole)",
                "Pitch-to-coverage attribution (unique workflow win)",
            ],
            notable_gaps=[
                "No dedicated mobile app",
                "Less “deep social intelligence” vs Brandwatch/Meltwater",
            ],
        ),
    ]

    add_title_slide(
        prs,
        title="Media Monitoring — Competitive Capabilities",
        subtitle="Top 3 products (Meltwater, Cision, Muck Rack) mapped into must-have / good-to-have / advantage features\nPrepared for AlmaLabs",
    )

    add_bullets_slide(
        prs,
        "Scope & goal",
        [
            "Focus: Media Monitoring (not full PR suite)",
            "Summarize core capabilities of top competitors",
            "Bucket features into: Must-have, Good-to-have, Competitive advantage opportunities",
        ],
        note="Sources: context.md; competitive-research/*.md (meltwater.md, cision.md, muck-rack.md, landscape-overview.md).",
    )

    add_competitor_snapshot_slide(prs, competitors)

    must_have_rows = [
        ("Online news monitoring at scale", ["✓", "✓", "✓"]),
        ("Advanced search (keywords + Boolean)", ["✓", "✓", "✓"]),
        ("Filters (date/source/geo/language)", ["✓", "✓", "✓"]),
        ("Real-time alerts (email) + digests", ["✓", "✓", "✓"]),
        ("Deduplication / story clustering", ["✓", "✓", "✓"]),
        ("Basic sentiment + volume trends", ["✓", "✓", "✓"]),
        ("Dashboards + scheduled reports", ["✓", "✓", "✓"]),
        ("Exports (CSV/PDF) + API/data feeds", ["✓", "✓", "✓ (add-on)"]),
        ("Licensed/premium content via partnerships", ["✓ (Factiva)", "✓ (Factiva/Nexis)", "✓ (LexisNexis)"]),
        ("Team collaboration (users/roles/sharing)", ["✓", "✓", "✓"]),
    ]
    add_matrix_slide(
        prs,
        "Must-have capabilities (table stakes)",
        must_have_rows,
        ["Meltwater", "Cision", "Muck Rack"],
        note="Interpretation: 'Must-have' = required to compete credibly in media monitoring for brands (news-first).",
    )

    good_to_have_rows = [
        ("Social listening (major platforms)", ["✓", "✓ (Brandwatch)", "Partial (Keyhole)"]),
        ("Mobile app w/ push alerts", ["✓", "✓", "—"]),
        ("Podcasts monitoring + transcription", ["✓", "✓", "✓"]),
        ("Broadcast TV/radio monitoring", ["✓ (regional)", "✓ (3K+ stations)", "✓ (TVEyes)"]),
        ("Print monitoring (clippings/OCR)", ["✓ (partners)", "✓ (strong in UK)", "✓ (LexisNexis)"]),
        ("Slack / Microsoft Teams integrations", ["Teams", "Slack+Teams", "Slack"]),
        ("Media database + outreach/pitching", ["✓", "✓ (largest)", "✓ (core strength)"]),
        ("Press release distribution", ["✓ (module)", "✓ (PR Newswire)", "✓ (GlobeNewswire)"]),
        ("Image/logo recognition", ["✓", "✓ (Brandwatch)", "—"]),
        ("Crisis scoring / risk signals", ["—", "✓ (React Score)", "—"]),
    ]
    add_matrix_slide(
        prs,
        "Good-to-have capabilities (common differentiators)",
        good_to_have_rows,
        ["Meltwater", "Cision", "Muck Rack"],
    )

    add_bullets_slide(
        prs,
        "Competitive advantage opportunities (what we can win on)",
        [
            "India-first licensed coverage (regional languages + tier-2/3 publishers) with transparent source list",
            "Accuracy layer: entity disambiguation + confidence scores + explainable matches",
            "Setup simplicity: templates + assisted query builder (reduce Boolean pain without AI hype)",
            "Fast alerts with clustering + executive summaries (email/Slack) and clear “why it matters”",
            "Compliance-by-design: licensing + snippet/click-through model + audit trail",
        ],
        note="CEO directive: keep positioning simple (“Media Monitoring”), avoid AI-buzzword marketing; use assistive UX where it lowers friction.",
    )

    add_bullets_slide(
        prs,
        "Where AlmaConnect News is today (baseline)",
        [
            "We already ingest large-scale public news via RSS + scraping",
            "We match against large prospect lists and deliver alerts (snippet + URL)",
            "Key gaps vs market: licensing, dedup/clustering, sentiment/topics, analytics/reporting, identity resolution",
        ],
        note="Source: product.md (current pipeline + limitations).",
    )

    add_bullets_slide(
        prs,
        "MVP checklist (news-first brand monitoring)",
        [
            "Licensing/partnerships for premium & reliable feeds (non-negotiable)",
            "Deduplication + story clustering",
            "Entity resolution + match confidence scoring",
            "Dashboards + scheduled reports (PDF/PPT-ready) + exports",
            "Slack/Teams + webhooks for alert routing",
        ],
    )

    add_bullets_slide(
        prs,
        "Next steps",
        [
            "Confirm target segment (India-first SMB vs global mid-market)",
            "Pick 1–2 licensing paths to validate cost + coverage",
            "Define MVP scope and timeline; run 3–5 customer discovery calls with demo prototype",
        ],
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate media monitoring competitive capabilities PPTX.")
    parser.add_argument(
        "--out",
        default="deliverables/media_monitoring_competitive_capabilities.pptx",
        help="Output .pptx path (default: deliverables/media_monitoring_competitive_capabilities.pptx)",
    )
    args = parser.parse_args()
    build_deck(Path(args.out))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
