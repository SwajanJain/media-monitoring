#!/usr/bin/env python3
"""
Generate AlmaLabs Media Monitoring Product Strategy PPTX
Self-explanatory presentation for CEO — reads without a presenter.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Design System ──────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x77, 0xB6)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)
YELLOW_ACCENT = RGBColor(0xF3, 0x9C, 0x12)
TABLE_HEADER_BG = NAVY
TABLE_ALT_BG = RGBColor(0xF8, 0xF9, 0xFA)

FONT_NAME = "Calibri"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_slide(prs):
    """Add a blank slide."""
    layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(layout)


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_size=14, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME, spacing_after=Pt(4)):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = spacing_after
    return p


def add_paragraph(tf, text, font_size=14, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, spacing_after=Pt(4), spacing_before=Pt(0)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.alignment = alignment
    p.space_after = spacing_after
    p.space_before = spacing_before
    return p


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a navy title bar at the top of the slide."""
    # Title background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # Title text
    tb = add_textbox(slide, Inches(0.75), Inches(0.15), Inches(11.5), Inches(0.55))
    set_text(tb.text_frame, title_text, font_size=26, bold=True, color=WHITE)

    if subtitle_text:
        tb2 = add_textbox(slide, Inches(0.75), Inches(0.65), Inches(11.5), Inches(0.35))
        set_text(tb2.text_frame, subtitle_text, font_size=13, color=RGBColor(0xBB, 0xCC, 0xDD))

    # Accent line
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), SLIDE_WIDTH, Inches(0.04)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = ACCENT
    shape2.line.fill.background()


def add_section_divider(slide, section_title, section_number=None):
    """Full-slide section divider."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    if section_number:
        tb = add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(0.8))
        set_text(tb.text_frame, section_number, font_size=18, color=ACCENT, bold=True)

    tb2 = add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.5))
    set_text(tb2.text_frame, section_title, font_size=40, bold=True, color=WHITE)


def add_body_text(slide, lines, left=0.75, top=1.4, width=11.8, font_size=14, line_spacing=Pt(6)):
    """Add multiple lines of body text."""
    tb = add_textbox(slide, Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        text = line.get("text", "")
        bold = line.get("bold", False)
        size = line.get("size", font_size)
        color = line.get("color", DARK_GRAY)
        indent = line.get("indent", 0)
        spacing_before = line.get("spacing_before", Pt(0))
        spacing_after = line.get("spacing_after", line_spacing)

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = spacing_after
        p.space_before = spacing_before
        if indent:
            p.level = indent
            p.indent = Inches(0.3 * indent)
    return tb


def add_table(slide, data, left=0.75, top=1.5, width=11.8, col_widths=None, font_size=10):
    """Add a formatted table. data[0] is header row."""
    rows = len(data)
    cols = len(data[0])
    table_height = Inches(0.35 * rows)

    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), table_height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = str(cell_text)
            p.font.size = Pt(font_size)
            p.font.name = FONT_NAME
            p.space_after = Pt(2)
            p.space_before = Pt(2)

            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            else:
                p.font.color.rgb = DARK_GRAY
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TABLE_ALT_BG
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def add_insight_box(slide, text, top=6.0, left=0.75, width=11.8):
    """Add an insight callout box."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    set_text(tf, "INSIGHT", font_size=9, bold=True, color=ACCENT)
    add_paragraph(tf, text, font_size=11, color=DARK_GRAY)
    return shape


def add_feature_block(slide, number, title, description, top, left=0.75, width=5.5):
    """Add a feature card."""
    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(0.35), Inches(0.35)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.margin_left = Inches(0)
    ctf.margin_right = Inches(0)
    ctf.margin_top = Inches(0)
    ctf.margin_bottom = Inches(0)
    set_text(ctf, str(number), font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Title
    tb = add_textbox(slide, Inches(left + 0.45), Inches(top - 0.02), Inches(width - 0.5), Inches(0.35))
    set_text(tb.text_frame, title, font_size=15, bold=True, color=NAVY)

    # Description
    tb2 = add_textbox(slide, Inches(left + 0.45), Inches(top + 0.35), Inches(width - 0.5), Inches(1.2))
    tf = tb2.text_frame
    tf.word_wrap = True
    set_text(tf, description, font_size=11, color=MED_GRAY, spacing_after=Pt(2))
    return tb2


# ── SLIDE BUILDERS ─────────────────────────────────────────────────────────────

def slide_01_title(prs):
    slide = add_slide(prs)
    # Navy background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.3), Inches(1.5), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Title
    tb = add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2))
    set_text(tb.text_frame, "Media Monitoring", font_size=44, bold=True, color=WHITE)
    add_paragraph(tb.text_frame, "Competitor Capabilities & Product Gaps", font_size=28, color=RGBColor(0xAA, 0xBB, 0xCC))

    # Subtitle
    tb2 = add_textbox(slide, Inches(1), Inches(4.2), Inches(10), Inches(0.5))
    set_text(tb2.text_frame, "What the top players offer, where they fail, and what Alma News should build", font_size=16, color=RGBColor(0x88, 0x99, 0xAA))

    # Research base
    tb3 = add_textbox(slide, Inches(1), Inches(5.2), Inches(10), Inches(1.2))
    tf = tb3.text_frame
    tf.word_wrap = True
    set_text(tf, "Research base:", font_size=12, bold=True, color=ACCENT)
    add_paragraph(tf, "4 competitor product teardowns (Meltwater, Cision, Muck Rack, Brand24)", font_size=12, color=RGBColor(0x99, 0xAA, 0xBB))
    add_paragraph(tf, "425 Reddit threads from r/PublicRelations with full comment trees", font_size=12, color=RGBColor(0x99, 0xAA, 0xBB))
    add_paragraph(tf, "Feature-by-feature comparison across 13 dimensions", font_size=12, color=RGBColor(0x99, 0xAA, 0xBB))

    # Footer
    tb4 = add_textbox(slide, Inches(1), Inches(6.8), Inches(10), Inches(0.4))
    set_text(tb4.text_frame, "February 2026  |  Swajan  |  AlmaLabs", font_size=11, color=RGBColor(0x77, 0x88, 0x99))


def slide_02_exec_summary(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "Executive Summary")

    points = [
        {"label": "The market", "text": "We researched the top 3 media monitoring companies \u2014 Meltwater ($500M revenue), Cision ($850-950M), Muck Rack ($143M). All charge $15K\u2013$100K/year. Users are unhappy with all three. G2 ratings: 3.8\u20134.6. Across 425 Reddit threads in the largest PR community (45K members), the common message is that none of these tools work well."},
        {"label": "The problem", "text": "They collect articles but don't help users understand them. Sentiment analysis is inaccurate (65\u201375% accuracy). Reports need to be rebuilt manually. Users don't trust the output."},
        {"label": "Where we are", "text": "Alma News already has the data side working \u2014 monitoring, scraping, matching, clustering, sentiment \u2014 for global clients. The missing pieces are on the output side: dashboard, alerts, reports, and smarter summarization."},
        {"label": "This deck", "text": "Maps every competitor feature into three categories \u2014 must have, good to have, and competitive advantage \u2014 and identifies 3 structural advantages that competitors can't copy because doing so would hurt their sales process, pricing model, or existing client workflows."},
    ]

    top = 1.5
    for i, pt in enumerate(points):
        # Number
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.75), Inches(top), Inches(0.4), Inches(0.4)
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = ACCENT
        num_shape.line.fill.background()
        ntf = num_shape.text_frame
        ntf.margin_left = ntf.margin_right = ntf.margin_top = ntf.margin_bottom = Inches(0)
        set_text(ntf, str(i + 1), font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        # Label
        tb = add_textbox(slide, Inches(1.3), Inches(top - 0.02), Inches(1.5), Inches(0.35))
        set_text(tb.text_frame, pt["label"].upper(), font_size=10, bold=True, color=ACCENT)

        # Text
        tb2 = add_textbox(slide, Inches(1.3), Inches(top + 0.3), Inches(11), Inches(0.9))
        set_text(tb2.text_frame, pt["text"], font_size=13, color=DARK_GRAY, spacing_after=Pt(2))

        top += 1.4


def slide_03_who_they_are(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "The Global Market \u2014 Who They Are",
                  "Cision makes 2x Meltwater's revenue with half the G2 rating \u2014 they're locked in by contracts, not product quality")

    data = [
        ["Metric", "Meltwater", "Cision", "Muck Rack", "Brand24"],
        ["Positioning", "Media intelligence\nBroadest data, AI-first", "End-to-end comms\nLargest DB, owns PR Newswire", "PR management\nBest UX, workflow-focused", "Affordable monitoring\nSelf-serve, SMB-focused"],
        ["Revenue", "~$500M", "~$850-950M", "$143.1M", "~$12M ARR"],
        ["Customers", "27,000+", "100,000+", "~5,000-6,000", "~3,685"],
        ["ARPU (annual)", "~$18,500", "~$15-25K", "$23,850", "$3,168"],
        ["G2 Rating", "4.0/5", "3.8/5", "4.6/5", "4.6/5"],
        ["G2 Ease of Use", "7.8/10", "7.2/10", "9.0/10", "9.2/10"],
        ["G2 Support", "Medium-High", 'Low ("black hole")', "9.3/10", "9.5/10"],
        ["Target Segment", "Enterprise, Mid-Market", "Enterprise, Mid-Market", "Mid-Market \u2192 Enterprise", "SMB"],
    ]
    add_table(slide, data, top=1.4, col_widths=[1.8, 2.5, 2.5, 2.5, 2.5], font_size=10)

    add_insight_box(slide,
        "Muck Rack and Brand24 both score 4.6 on G2 but Muck Rack charges 7x more. The difference is the sales process, not what users experience.",
        top=6.2)


def slide_04_how_they_sell(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "The Global Market \u2014 How They Sell & What Breaks",
                  "Brand24 is built for users to try before buying. The other three are built for sales teams to close before users try.")

    data = [
        ["Metric", "Meltwater", "Cision", "Muck Rack", "Brand24"],
        ["Sales model", "Sales-led, demo required", "Sales-led, demo required", "Sales-led, demo required", "Self-serve"],
        ["Onboarding", "2-4 weeks", "2-4 weeks", "1-2 days", "Under 10 minutes"],
        ["Pricing", "$15K-100K+/yr", "$6K-100K+/yr", "$5K-53K+/yr", "$948-4,788/yr"],
        ["Public pricing", "No", "No", "No", "Yes"],
        ["Monthly billing", "No", "No", "No", "Yes"],
        ["Free trial", "No", "No", "No", "Yes (14 days)"],
        ["Renewal hikes", "5-50%", "20-50%", "17% (new)", "Not reported"],
        ["Auto-renewal traps", "30-60 day window", "Buried on page 47", "Email CSM + legal", "Not reported"],
        ["Sentiment accuracy", "Won't disclose", "65-75%", "70-75% (rules)", "95% (EN), 60-70% others"],
    ]
    add_table(slide, data, top=1.4, col_widths=[1.8, 2.5, 2.5, 2.5, 2.5], font_size=10)


def slide_05_generic_flow(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "How Media Monitoring Works \u2014 The Generic Flow",
                  "Every tool follows the same 5 steps. Where they differ is how much manual work the client still does.")

    steps = [
        ("1", "Setup", "Client tells the tool what to track\n(brand name, competitors, keywords)", ACCENT),
        ("2", "Collection", "Tool pulls articles from news,\nsocial, blogs, wire services", ACCENT),
        ("3", "Processing", "Filter irrelevant results, group\nrelated articles, tag sentiment", ACCENT),
        ("4", "Presentation", "Client sees results in dashboard\n\u2014 articles, charts, scores", ACCENT),
        ("5", "Delivery", "Alerts (email, Slack) and\nreports (PDF, CSV)", ACCENT),
    ]

    x_start = 0.5
    step_width = 2.3
    arrow_width = 0.25
    y_top = 2.0

    for i, (num, title, desc, color) in enumerate(steps):
        x = x_start + i * (step_width + arrow_width + 0.1)

        # Step box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y_top), Inches(step_width), Inches(1.8)
        )
        box.fill.solid()
        if i < 3:
            box.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
            border_color = ACCENT
        else:
            box.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
            border_color = RED_ACCENT
        box.line.color.rgb = border_color
        box.line.width = Pt(1.5)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)
        set_text(tf, f"Step {num}", font_size=10, bold=True, color=border_color)
        add_paragraph(tf, title, font_size=16, bold=True, color=NAVY, spacing_before=Pt(4))
        add_paragraph(tf, desc, font_size=10, color=MED_GRAY, spacing_before=Pt(6))

        # Arrow
        if i < 4:
            arrow_x = x + step_width + Inches(0.03)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x), Inches(y_top + 0.7), Inches(arrow_width), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            arrow.line.fill.background()

    # Alma News status
    tb = add_textbox(slide, Inches(0.75), Inches(4.5), Inches(11.8), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, "Alma News today covers Steps 1-3:", font_size=14, bold=True, color=NAVY)
    add_paragraph(tf, "Keyword setup, article scraping, matching, clustering, sentiment. We already run 11-category sentiment analysis on our alumni product \u2014 adapting for media monitoring is proven capability, not new R&D.", font_size=13, color=DARK_GRAY)
    add_paragraph(tf, "The gaps are in Steps 4-5: presentation and delivery.", font_size=14, bold=True, color=RED_ACCENT, spacing_before=Pt(8))


def slide_product_overview(prs, company, positioning, user_flow, unique_features, main_problem, problem_color=RED_ACCENT):
    """Generic product overview slide."""
    slide = add_slide(prs)
    add_title_bar(slide, f"{company} \u2014 Product Overview", positioning)

    # Left column: User flow
    tb = add_textbox(slide, Inches(0.75), Inches(1.3), Inches(0.5), Inches(0.3))
    set_text(tb.text_frame, "HOW IT WORKS", font_size=9, bold=True, color=ACCENT)

    top = 1.65
    for step_title, step_desc in user_flow:
        tb = add_textbox(slide, Inches(0.75), Inches(top), Inches(6.5), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        set_text(tf, step_title, font_size=12, bold=True, color=NAVY, spacing_after=Pt(1))
        add_paragraph(tf, step_desc, font_size=10, color=MED_GRAY, spacing_after=Pt(1))
        top += 0.55

    # Right column: Unique features
    right_left = 7.8
    tb2 = add_textbox(slide, Inches(right_left), Inches(1.3), Inches(5), Inches(0.3))
    set_text(tb2.text_frame, "UNIQUE FEATURES", font_size=9, bold=True, color=ACCENT)

    tb3 = add_textbox(slide, Inches(right_left), Inches(1.65), Inches(5), Inches(3.5))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    first = True
    for feat in unique_features:
        if first:
            set_text(tf3, f"\u2022  {feat}", font_size=11, color=DARK_GRAY, spacing_after=Pt(6))
            first = False
        else:
            add_paragraph(tf3, f"\u2022  {feat}", font_size=11, color=DARK_GRAY, spacing_after=Pt(6))

    # Main problem box at bottom
    prob_top = max(top + 0.2, 5.2)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(prob_top), Inches(11.8), Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
    shape.line.color.rgb = RED_ACCENT
    shape.line.width = Pt(1)

    tf4 = shape.text_frame
    tf4.word_wrap = True
    tf4.margin_left = Inches(0.15)
    tf4.margin_top = Inches(0.08)
    set_text(tf4, "THE MAIN PROBLEM", font_size=9, bold=True, color=RED_ACCENT)
    add_paragraph(tf4, main_problem, font_size=12, color=DARK_GRAY)


def slide_06_meltwater(prs):
    slide_product_overview(prs,
        company="Meltwater",
        positioning="Media intelligence \u2014 broadest data coverage, AI-first analytics  |  ~$500M revenue  |  27,000+ customers",
        user_flow=[
            ("Setup", "Boolean queries or AI Search Assistant (plain English \u2192 Boolean). 7+ wizard steps. 2-4 weeks with CSM."),
            ("Monitor", "\"Explore\" feed shows articles with headline, source, sentiment, reach. Also \"Monitor\" for saved feeds \u2014 two places for similar things."),
            ("Analyze", "Customizable widgets \u2014 sentiment trends, volume, geographic breakdown, share of voice. City-level data (unique)."),
            ("Report", "AI executive summaries + cover images. PDF, CSV, JSON feeds. Shareable live links."),
            ("Alerts", "Real-time email + mobile push. AI explains WHY a spike happened. Slack & Teams."),
        ],
        unique_features=[
            "GenAI Lens \u2014 tracks brands in ChatGPT/Gemini/Claude answers",
            "AI Search Assistant \u2014 plain English to Boolean",
            "Entity-level sentiment \u2014 brand-specific even in multi-brand articles",
            "Image/logo recognition in social posts",
            "270,000+ news sources, 500M+ pieces/day ingested",
        ],
        main_problem="10+ sections with overlapping functionality. Powerful but takes weeks to learn. Users: \"the tool that can do everything but takes forever to figure out.\""
    )


def slide_07_cision(prs):
    slide_product_overview(prs,
        company="Cision",
        positioning="End-to-end communications \u2014 largest database, owns PR Newswire  |  ~$850-950M revenue  |  100,000+ customers",
        user_flow=[
            ("Setup", "Boolean queries manually \u2014 no AI assist. Must learn AND/OR/NOT/NEAR. 2-4 weeks, 3 training sessions."),
            ("Monitor", "\"Mention Streams\" inbox-style list. Headline, sentiment badge, React Score, engagement. Sort by date/relevance/reach."),
            ("Discover", "Largest journalist database (700K-1M+ contacts). Search by beat, geography, outlet, influence."),
            ("Distribute", "PR Newswire (owned) \u2014 write, target, distribute press releases, then track pickup. Unique closed loop."),
            ("Analyze", "Custom dashboard widgets. \"React Score\" \u2014 crisis detection from 17 NLP models, 0-600 scale. Unique."),
        ],
        unique_features=[
            "PR Newswire (owned) \u2014 distribute + track pickup in one system",
            "React Score \u2014 17 NLP models for crisis detection",
            "Largest journalist database (700K-1M+ contacts)",
            "300,000+ editorial calendar opportunities",
            "Human-curated print monitoring in UK (95-99% accuracy)",
        ],
        main_problem="Feels like 4 products stitched together \u2014 Cision + Brandwatch + TrendKite + PR Newswire (all acquisitions). Brandwatch requires a different login. G2 ease of use: 7.2/10 \u2014 worst in category."
    )


def slide_08_muckrack(prs):
    slide_product_overview(prs,
        company="Muck Rack",
        positioning="PR management \u2014 journalist-first, best UX, workflow-focused  |  $143.1M revenue  |  G2 ease of use: 9.0/10",
        user_flow=[
            ("Setup", "\"Trackers\" \u2014 Boolean searches across 5 sources simultaneously. Deduplicated. 1-2 days with CSM."),
            ("Monitor", "Article feed + \"Spike Alerts\" (ML detects unusual volume, AI summarizes why). Topic filter auto-excludes noise."),
            ("Discover", "~250K journalist contacts. Journalists maintain their own profiles for free hosting. Self-updating."),
            ("Pitch", "Integrated email with tracking. \"Pitch Coverage Detection\" \u2014 auto-detects when pitched journalist publishes."),
            ("Report", "Coverage Reports auto-populate. \"PR Hit Score\" \u2014 composite metric with adjustable weights."),
            ("AI", "\"Generative Pulse\" \u2014 tracks brand in ChatGPT/Claude/Gemini/Perplexity. $3K extra/year."),
        ],
        unique_features=[
            "Two-sided marketplace \u2014 journalists maintain own profiles",
            "Pitch Coverage Detection \u2014 links pitches to published coverage",
            "Media List Agent (beta Feb 2026) \u2014 first AI agent in PR SaaS",
            "Single codebase \u2014 no acquisitions, everything feels consistent",
            "PR Hit Score \u2014 customizable composite metric",
        ],
        main_problem="US-centric \u2014 first non-US office just opened (UK, Oct 2024). No mobile app. Pricing rising \u2014 $34K/year for 12 seats, 17% renewal hike. AI acknowledged as lagging vs Meltwater."
    )


def slide_09_brand24(prs):
    slide_product_overview(prs,
        company="Brand24",
        positioning="Affordable monitoring \u2014 self-serve, transparent pricing, SMB-focused  |  ~$12M ARR  |  G2: 4.6/5",
        user_flow=[
            ("Setup", "Sign up online. Type keyword, select language, click \"Create Project.\" Results in minutes. Under 10 min total."),
            ("Monitor", "Mentions feed from news, social, blogs, Reddit. Color-coded sentiment. Update frequency by plan tier."),
            ("Analyze", "Pre-built charts \u2014 volume, sentiment, sources, geographic, emoji analysis. No custom dashboards."),
            ("AI", "\"AI Brand Assistant\" \u2014 chat interface. Ask questions about your data, get answers with cited mentions."),
            ("LLM Listening", "\"Chatbeat\" \u2014 monitors 8+ AI models (ChatGPT, Gemini, Claude, Perplexity, DeepSeek, Grok, Copilot)."),
            ("Reports", "One-click infographic reports. White-labeled PDFs from $129/month plan."),
        ],
        unique_features=[
            "Self-serve onboarding \u2014 under 10 min, no human involved",
            "Public pricing ($79-$399/month), monthly billing",
            "14-day free trial, no credit card required",
            "Chatbeat \u2014 monitors 8+ LLMs, most comprehensive AI citation tracking",
            "NeurIPS-published sentiment research",
            "6-emotion detection (Admiration, Joy, Anger, Disgust, Fear, Sadness)",
        ],
        main_problem="No journalist database. No paywalled content (no WSJ/FT/Bloomberg). Standard Twitter API misses 70-90% of tweets. Non-English sentiment: 60-70%. Being absorbed into Adobe through Semrush \u2014 future uncertain."
    )


def slide_10_the_gap(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "The One Gap That Matters", "Data collection gaps are closing. The output gap stays wide open.")

    # Context line
    tb = add_textbox(slide, Inches(0.75), Inches(1.4), Inches(11.8), Inches(0.6))
    set_text(tb.text_frame,
        "Competitors are adding AI features fast \u2014 Meltwater has AI summaries, Brand24 claims 95% English sentiment, 3 of 4 now track AI citations.",
        font_size=13, color=MED_GRAY)

    # Big statement
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.2), Inches(11.8), Inches(0.7)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    stf = shape.text_frame
    stf.margin_left = Inches(0.2)
    stf.margin_top = Inches(0.1)
    set_text(stf, "The output is broken.", font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Evidence
    lines = [
        {"text": "Every monitoring tool ends the same way: a dashboard with \"327 mentions, 68% positive, 80M potential reach.\"", "size": 13},
        {"text": "No PR team sends this to leadership. They export, open PowerPoint, and spend hours rebuilding.", "size": 13},
        {"text": "", "size": 8},
        {"text": "This is the #1 complaint category:", "size": 13, "bold": True, "color": NAVY},
        {"text": "\u2022  30% of Meltwater's negative G2 reviews mention reporting", "size": 13},
        {"text": "\u2022  30% of Muck Rack's negative G2 reviews mention reporting", "size": 13},
        {"text": "\u2022  Reddit users describe the same cycle: \"pull data, open PowerPoint, rebuild\"", "size": 13},
        {"text": "", "size": 8},
        {"text": "Competitors can't fix this \u2014 changing the output format breaks templates that 27,000+ clients (Meltwater) and 100,000+ clients (Cision) depend on.", "size": 12, "color": MED_GRAY},
    ]
    add_body_text(slide, lines, top=3.1, font_size=13)

    # Example briefing box
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(5.6), Inches(11.8), Inches(1.5)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0xF0, 0xF7, 0xEB)
    shape2.line.color.rgb = GREEN_ACCENT
    shape2.line.width = Pt(1)

    etf = shape2.text_frame
    etf.word_wrap = True
    etf.margin_left = Inches(0.15)
    etf.margin_top = Inches(0.08)
    set_text(etf, "WHAT THE OUTPUT SHOULD LOOK LIKE", font_size=9, bold=True, color=GREEN_ACCENT)
    add_paragraph(etf,
        "\"Coverage this week focused on the CFO's comments at Davos. 14 outlets picked it up, led by Reuters. Tone was skeptical \u2014 most articles questioned the revenue guidance. This is a shift from last month. The Reuters piece is being cited by ChatGPT. Recommended: prepare a response before earnings.\"",
        font_size=12, color=DARK_GRAY, spacing_before=Pt(4))
    add_paragraph(etf, "That's one LLM call on top of clustered articles. We already cluster. The gap is the last mile.", font_size=11, bold=True, color=NAVY, spacing_before=Pt(6))


def slide_11_must_have_core(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "Must Have \u2014 The Core View",
                  "8 features to complete the PR team's workflow. Without any one, the product doesn't work.")

    # 3 modes
    tb = add_textbox(slide, Inches(0.75), Inches(1.35), Inches(11.8), Inches(0.3))
    set_text(tb.text_frame, "A PR person's day has 3 modes:  Morning check  |  Active monitoring  |  Weekly reporting", font_size=12, bold=True, color=NAVY)

    # Feature 1: Monitoring Feed
    add_feature_block(slide, "1", "Monitoring Feed \u2014 the screen where the user lives",
        "Not a list of articles. A list of stories. Clusters as cards \u2014 each shows: AI-generated headline, article count, top sources (logos), sentiment badge, trend arrow (\u2191 growing / \u2192 steady / \u2193 fading), time range. Click to expand. Different from every competitor: stories first, articles second.",
        top=1.9, left=0.75, width=5.5)

    # Feature 2: Boolean Search
    add_feature_block(slide, "2", "Boolean Search + Filters",
        "AND, OR, NOT, NEAR operators. Exact phrases. Parentheses for grouping. Sidebar filters: date range, source type, sentiment, geography, outlet name. Without Boolean, agencies can't set up precise monitoring. Separates a toy from a tool.",
        top=1.9, left=7.0, width=5.5)

    # Feature 3: Deduplication
    add_feature_block(slide, "3", "Deduplication \u2014 invisible but critical",
        "User never sees this as a feature. When Reuters publishes something and 20 outlets syndicate it, the feed shows one card: \"Reuters: Q4 results \u2014 appeared in 20 outlets.\" Without this, the feed is noise. Must ship with source expansion.",
        top=4.4, left=0.75, width=11.5)

    add_insight_box(slide,
        "We show stories first, articles second. Meltwater shows articles across 10+ modules. Cision shows mention streams. Simpler, faster for the morning check.",
        top=6.1)


def slide_12_must_have_intel(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "Must Have \u2014 Intelligence Layer",
                  "What the system tells you without being asked")

    # Feature 4: Sentiment
    add_feature_block(slide, "4", "Sentiment \u2014 with a one-line reason",
        "Badge alone is useless. We show badge + why.\nNegative (red): \"Questions product reliability after customer complaints\"\nPositive (green): \"Highlights new feature launch, quotes CEO favorably\"\nOne LLM call per article. We already run 11-category sentiment on our alumni product \u2014 this is adapting proven capability, not new R&D.",
        top=1.5, left=0.75, width=5.5)

    # Feature 5: Email Alerts
    add_feature_block(slide, "5", "Email Alerts \u2014 two types, nothing more",
        "Daily digest: \"Yesterday: 23 articles, 4 stories. Top story: [X]. One negative cluster growing.\"\nSpike alert: \"12 articles in 2 hours about [X]. Normal: 3-4/day.\"\nTwo modes only: daily summary + \"something is happening right now.\"",
        top=1.5, left=7.0, width=5.5)

    # Feature 6: Volume Spikes
    add_feature_block(slide, "6", "Volume Spikes \u2014 visual indicator on the feed",
        "When a story gains traction faster than normal, the cluster card gets a visual flag. \"This story: 3 articles yesterday, 14 today, still growing.\" The spike email handles it when you're away. The visual indicator handles it when you're in the product.",
        top=4.2, left=0.75, width=11.5)

    add_insight_box(slide,
        "The difference between a tool that shows what happened and one that tells you what's happening right now. PR teams need to catch stories before they blow up, not after.",
        top=6.1)


def slide_13_must_have_output(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "Must Have \u2014 Output Layer",
                  "What you take away from the product")

    # Feature 7: Dashboard
    add_feature_block(slide, "7", 'Dashboard \u2014 the "screenshot for leadership" view',
        "Separate tab. Four charts: volume over time (line), sentiment breakdown (bar), top sources, top stories. Time selector: 7d / 30d / 90d / custom. Must look clean enough to screenshot and paste into a slide \u2014 that's literally how people use dashboards at Meltwater today.",
        top=1.5, left=0.75, width=5.5)

    # Feature 8: Export
    add_feature_block(slide, "8", "Export \u2014 PDF and CSV",
        "PDF: One-click dashboard export. Charts + top stories + sentiment. Clean enough to email.\nCSV: Raw data \u2014 every article with date, source, headline, URL, sentiment, cluster ID. For Excel teams.\nIn V2, this becomes the AI briefing \u2014 auto-generated narrative report.",
        top=1.5, left=7.0, width=5.5)

    # What's NOT in V1
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(4.2), Inches(11.8), Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    set_text(tf, "NOT IN V1", font_size=9, bold=True, color=MED_GRAY)
    add_paragraph(tf, "AI briefing report  \u2022  Slack/Teams  \u2022  API access  \u2022  Competitive benchmarking  \u2022  Mobile app  \u2022  AI citation tracking", font_size=12, color=MED_GRAY)
    add_paragraph(tf, "These move to Good to Have and Competitive Advantage.", font_size=11, color=MED_GRAY)


def slide_14_good_to_have_quick(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "Good to Have \u2014 Quick Wins",
                  "Weeks to build, not months. The natural V2 requests after using V1 for a month.")

    add_feature_block(slide, "1", "Slack/Teams Alerts",
        "PR teams live in Slack. Spike alert in Slack vs email = 30 seconds vs 30 minutes. Same two message types as email \u2014 morning digest + spike. Click \u2192 opens cluster in product. Webhook integration. Meltwater/Cision have both. Muck Rack: Slack only.",
        top=1.5, left=0.75, width=5.5)

    add_feature_block(slide, "2", "API Access",
        "Agencies want raw data in Tableau, PowerBI, Google Sheets. REST API with article/cluster/sentiment endpoints. API key per client, rate-limited by tier. Premium upsell \u2014 Muck Rack gates behind $25K+ Premier plan.",
        top=1.5, left=7.0, width=5.5)

    add_feature_block(slide, "3", "Mobile Access (PWA)",
        "Not a native app \u2014 progressive web app. Same product, responsive layout. Cluster cards stack vertically, charts resize. No app store. 80% of native app value at 10% of the cost. Meltwater/Cision have native apps. Muck Rack has nothing.",
        top=4.0, left=0.75, width=11.5)


def slide_15_good_to_have_strategic(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "Good to Have \u2014 Strategic Decisions",
                  "Not just features \u2014 each one is a product direction that requires a real decision.")

    add_feature_block(slide, "4", "Competitive Benchmarking + Share of Voice",
        "Second most common question: \"How are we doing vs [competitor]?\" User sets up monitoring for own brand + competitors. \"Compare\" tab: article volume side-by-side, sentiment comparison, source overlap, share of voice %. This is what agencies pay for \u2014 Meltwater charges enterprise pricing partly for this.",
        top=1.5, left=0.75, width=11.5)

    add_feature_block(slide, "5", "Multi-Language Sentiment",
        "Global clients get non-English coverage. Brand24 drops to 60-70% for non-English. Muck Rack: English-only NLTK. LLM-based sentiment (already built for one-line reasons in V1) handles multiple languages better \u2014 GPT-4/Claude strong across European and Asian languages. Needs testing per language before selling.",
        top=3.5, left=0.75, width=11.5)


def slide_16_competitive_advantage(prs):
    """Split into 2 slides for readability."""
    # Slide 16a: AI Briefing (the big one)
    slide = add_slide(prs)
    add_title_bar(slide, "Competitive Advantage \u2014 AI Briefing That Replaces the Report",
                  "What competitors can't copy because doing so would break how they make money")

    lines = [
        {"text": "Every PR team using Meltwater, Cision, or Muck Rack does the same thing every Monday:", "size": 14, "color": DARK_GRAY},
        {"text": "Export data from dashboard \u2192 Open PowerPoint \u2192 Spend hours rebuilding into something leadership can read", "size": 14, "bold": True, "color": NAVY},
        {"text": "30% of negative G2 reviews for both Meltwater and Muck Rack mention this exact problem.", "size": 13, "color": MED_GRAY, "spacing_before": Pt(8)},
    ]
    add_body_text(slide, lines, top=1.4)

    # Example briefing
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(3.0), Inches(11.8), Inches(1.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF7, 0xEB)
    shape.line.color.rgb = GREEN_ACCENT
    shape.line.width = Pt(1.5)

    etf = shape.text_frame
    etf.word_wrap = True
    etf.margin_left = Inches(0.2)
    etf.margin_top = Inches(0.1)
    set_text(etf, "WHAT THE CLIENT RECEIVES EVERY MORNING", font_size=9, bold=True, color=GREEN_ACCENT)
    add_paragraph(etf,
        "\"Coverage this week focused on the CFO's comments at Davos. 14 outlets picked it up, led by Reuters. Tone was skeptical \u2014 most articles questioned the revenue guidance. This is a shift from last month, when coverage was neutral. The Reuters piece is the most shared. Recommended: prepare a response to the revenue narrative before earnings next week.\"",
        font_size=13, color=DARK_GRAY, spacing_before=Pt(6))

    # How + Why
    tb = add_textbox(slide, Inches(0.75), Inches(4.9), Inches(5.5), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, "HOW IT WORKS", font_size=10, bold=True, color=ACCENT)
    add_paragraph(tf, "We already cluster articles. Run an LLM on each day's clusters \u2014 what happened, why it matters, what changed from yesterday. One API call per cluster. Delivered as email or in-app summary.", font_size=12, color=DARK_GRAY, spacing_before=Pt(4))

    tb2 = add_textbox(slide, Inches(7.0), Inches(4.9), Inches(5.5), Inches(1.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "WHY COMPETITORS CAN'T DO THIS", font_size=10, bold=True, color=RED_ACCENT)
    add_paragraph(tf2, "Replacing dashboards with AI summaries reduces time-in-product, which Meltwater and Cision use to justify $15K-100K/year pricing. A briefing that answers the question in 30 seconds undermines their pricing model.", font_size=12, color=DARK_GRAY, spacing_before=Pt(4))

    # Slide 16b: Self-Serve + Reddit
    slide2 = add_slide(prs)
    add_title_bar(slide2, "Competitive Advantage \u2014 Business Model + Data",
                  "Structural advantages that can't be copied without breaking how competitors operate")

    # Self-serve
    add_feature_block(slide2, "2", "Self-Serve Onboarding + Transparent Pricing",
        "Meltwater: 2-4 weeks, CSM required. Cision: 2-4 weeks, 3 training sessions. Muck Rack: 1-2 days. All annual-only, 17-50% renewal hikes.\n\nOur model: Website \u2192 see pricing \u2192 pick plan \u2192 enter keywords \u2192 see results in minutes. Monthly billing. Cancel anytime.\n\nBrand24 built $12M ARR on this exact model. The big 3 can't switch \u2014 their revenue depends on sales reps closing annual contracts. Self-serve kills the sales team. Monthly billing breaks revenue recognition.",
        top=1.5, left=0.75, width=11.5)

    # Reddit
    add_feature_block(slide2, "3", "Reddit Deep Monitoring",
        "Reddit is where honest brand opinions live and crises surface first. But tools either ignore it or scan headlines only.\n\nMeltwater: surface level. Cision: via Brandwatch (separate login). Muck Rack: nothing. Brand24: standard API (misses 70-90%).\n\nOur approach: Same monitoring feed as news. \"Your brand in r/technology \u2014 3 threads, 47 comments, sentiment negative, growing.\" Click for full comment tree with per-comment sentiment. Reddit API + same matching/sentiment pipeline.",
        top=4.0, left=0.75, width=11.5)

    add_insight_box(slide2,
        "These advantages are durable not because they're hard to build, but because they'd break competitors' business models. AI briefings reduce engagement. Self-serve kills sales teams. Reddit requires rethinking the data model. We have none of these constraints.",
        top=6.2)


def slide_17_pricing_pyramid(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "The Pricing Pyramid \u2014 Where the Gap Is",
                  "The $1K-15K segment has no credible product designed for it")

    segments = [
        {"label": "Ultra-Premium  $40K-100K+/yr", "tag": "OVERSATURATED", "tag_color": MED_GRAY,
         "players": "Meltwater, Cision", "buyers": "Fortune 500, top-10 agencies",
         "note": "Pay the most, still rebuild reports in PowerPoint.", "width": 4.0},
        {"label": "Premium  $15K-40K/yr", "tag": "CROWDED", "tag_color": MED_GRAY,
         "players": "Meltwater, Cision, Muck Rack", "buyers": "Mid-size companies, mid-size agencies",
         "note": "Muck Rack wins on UX. Others feel bloated. 2-4 weeks to onboard.", "width": 6.0},
        {"label": "Lower-Mid  $5K-15K/yr", "tag": "UNDERSERVED", "tag_color": RED_ACCENT,
         "players": "Muck Rack entry, Cision entry, Brand24 top", "buyers": "Small PR teams (2-5), small agencies",
         "note": "Worst version of each tool \u2014 limited seats, no API, least CSM attention.", "width": 8.5},
        {"label": "Budget  $1K-5K/yr", "tag": "ONE PLAYER, LEAVING", "tag_color": RED_ACCENT,
         "players": "Brand24 alone", "buyers": "Startups, solo consultants, SMBs",
         "note": "G2: 4.6 but being acquired by Adobe/Semrush \u2014 future uncertain.", "width": 10.5},
        {"label": "Below $1K/yr", "tag": "NOBODY CREDIBLE", "tag_color": RED_ACCENT,
         "players": "Google Alerts (free), Press Ranger ($49/mo)", "buyers": "Freelancers, tiny startups",
         "note": "Basic alerts only. No clustering, sentiment, or dashboards.", "width": 11.8},
    ]

    top = 1.4
    for seg in segments:
        left = (13.333 - seg["width"]) / 2
        # Row background
        row_height = 0.85
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(seg["width"]), Inches(row_height)
        )
        if seg["tag_color"] == RED_ACCENT:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
            shape.line.color.rgb = RED_ACCENT
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
            shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        shape.line.width = Pt(0.75)

        # Content
        tb = add_textbox(slide, Inches(left + 0.15), Inches(top + 0.05), Inches(seg["width"] - 0.3), Inches(row_height - 0.1))
        tf = tb.text_frame
        tf.word_wrap = True
        set_text(tf, f"{seg['label']}     {seg['tag']}", font_size=12, bold=True, color=NAVY, spacing_after=Pt(1))
        add_paragraph(tf, f"Players: {seg['players']}  |  Buyers: {seg['buyers']}", font_size=9, color=MED_GRAY, spacing_after=Pt(0))
        add_paragraph(tf, seg["note"], font_size=9, color=seg["tag_color"], spacing_after=Pt(0))

        top += row_height + 0.1

    # Gap callout
    gap_top = top + 0.15
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(gap_top), Inches(11.8), Inches(0.85)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = NAVY
    shape2.line.fill.background()
    gtf = shape2.text_frame
    gtf.word_wrap = True
    gtf.margin_left = Inches(0.2)
    gtf.margin_top = Inches(0.08)
    set_text(gtf, "THE GAP: $1K-15K", font_size=12, bold=True, color=ACCENT)
    add_paragraph(gtf, "Brand24 is the only credible product under $5K and they're being acquired. $5K-15K clients get stripped-down enterprise tools. Nobody designs a product for this segment. The opportunity: build the best product for the $1K-15K client from scratch.", font_size=12, color=WHITE, spacing_before=Pt(2))


def slide_18_summary(prs):
    slide = add_slide(prs)
    add_title_bar(slide, "Summary")

    points = [
        {
            "num": "1",
            "title": "Users are unhappy with every major player",
            "text": "Meltwater ($500M), Cision ($850-950M), Muck Rack ($143M) \u2014 G2 ratings 3.8-4.6. Across 425 Reddit threads: none of these tools work well. Users stay because there's no credible alternative."
        },
        {
            "num": "2",
            "title": "The one gap that stays open: the output is broken",
            "text": "Every tool collects articles. No tool produces a report leadership can read without rebuilding in PowerPoint. 30% of negative G2 reviews for Meltwater and Muck Rack are about reporting. Competitors can't fix this without breaking templates 27,000+ clients depend on."
        },
        {
            "num": "3",
            "title": "The $1K-15K segment has no credible product",
            "text": "Brand24 is the only option under $5K and they're being acquired. $5K-15K clients get stripped-down enterprise tools. Nobody designs a product for this segment."
        },
        {
            "num": "4",
            "title": "Three advantages competitors can't copy",
            "text": "Not because they're hard to build, but because they'd break how competitors make money. AI briefings reduce time-in-product (hurts engagement metrics). Self-serve kills the sales team (hurts revenue model). Reddit deep monitoring requires rethinking the data model. We don't have any of these constraints."
        },
    ]

    top = 1.6
    for pt in points:
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.75), Inches(top), Inches(0.4), Inches(0.4)
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = ACCENT
        num_shape.line.fill.background()
        ntf = num_shape.text_frame
        ntf.margin_left = ntf.margin_right = ntf.margin_top = ntf.margin_bottom = Inches(0)
        set_text(ntf, pt["num"], font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        tb = add_textbox(slide, Inches(1.3), Inches(top - 0.02), Inches(11), Inches(0.35))
        set_text(tb.text_frame, pt["title"], font_size=16, bold=True, color=NAVY)

        tb2 = add_textbox(slide, Inches(1.3), Inches(top + 0.35), Inches(11), Inches(0.8))
        set_text(tb2.text_frame, pt["text"], font_size=13, color=DARK_GRAY)

        top += 1.35


# ── MAIN ───────────────────────────────────────────────────────────────────────

def main():
    prs = new_presentation()

    # Slide 1: Title
    slide_01_title(prs)

    # Slide 2: Executive Summary
    slide_02_exec_summary(prs)

    # Section: Market
    s = add_slide(prs)
    add_section_divider(s, "The Global Market", "SECTION 1")

    # Slide 3-4
    slide_03_who_they_are(prs)
    slide_04_how_they_sell(prs)

    # Slide 5: Generic flow
    slide_05_generic_flow(prs)

    # Section: Competitors
    s = add_slide(prs)
    add_section_divider(s, "Competitor Deep Dives", "SECTION 2")

    # Slides 6-9: Product overviews
    slide_06_meltwater(prs)
    slide_07_cision(prs)
    slide_08_muckrack(prs)
    slide_09_brand24(prs)

    # Section: Gaps & Features
    s = add_slide(prs)
    add_section_divider(s, "The Gap & What to Build", "SECTION 3")

    # Slide 10: The One Gap
    slide_10_the_gap(prs)

    # Slides 11-13: Must Have
    slide_11_must_have_core(prs)
    slide_12_must_have_intel(prs)
    slide_13_must_have_output(prs)

    # Slides 14-15: Good to Have
    slide_14_good_to_have_quick(prs)
    slide_15_good_to_have_strategic(prs)

    # Slide 16: Competitive Advantage (2 slides)
    slide_16_competitive_advantage(prs)

    # Section: Market Position
    s = add_slide(prs)
    add_section_divider(s, "Market Position & Summary", "SECTION 4")

    # Slide 17: Pricing Pyramid
    slide_17_pricing_pyramid(prs)

    # Slide 18: Summary
    slide_18_summary(prs)

    # Save
    output_dir = "/Users/swajanjain/Documents/Projects/media-monitoring/deliverables"
    output_path = os.path.join(output_dir, "AlmaLabs_Media_Monitoring_Strategy.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
